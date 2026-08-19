"""
Document ingestion for per-client Knowledge Folders.

A client can be pointed at a folder of existing documents (SOWs,
discovery notes, requirements docs). This module extracts plain text
from each supported file, chunks it, embeds the chunks with the SAME
local model the transcript search uses (see core.embeddings), and
persists one .npz/.json sidecar pair per document so SearchService can
fold them into the same in-memory search matrix as session transcript
chunks.

LMA gap analysis 2026-08-07: three separate field reports trace back to
silent skips — a file that couldn't be processed just vanished from the
count, which reads identically to "nothing in this folder matters."
Every skip here is therefore COUNTED and carries an actionable reason
(missing optional dependency, encrypted/corrupt file, empty
extraction) rather than raising or disappearing. extract_text() never
raises for a bad *input* file; it only raises for genuine programming
errors.

Follow-up field data (2026-08-14): a real client Knowledge Folder — a
user's actual working corpus of SOWs, estimates, and RFP responses —
hit 87 skipped files out of what should have been the bulk of the
folder. Two distinct bugs, not one: SUPPORTED_EXTENSIONS claimed .pdf
and .docx support that pypdf/python-docx never shipped (advice to
`pip install` them is also unactionable inside a packaged app with a
bootstrapped venv, so that's gone too), and .xlsx/.pptx/.html had no
extractor at all despite openpyxl already being bundled for a different
feature. See the dispatch table below (_EXTRACTORS) and
NON_TEXT_EXTENSIONS, which now separate "we can't read this yet" from
"this was never going to be readable" (images, diagrams, media,
archives) — the field reports specifically called out the latter
reading like a defect when it isn't one.

Sidecar schema (one .npz + one .json per document, under
<recordings_dir>/doc_index/):
    {
        "model_id": str,             # see core.embeddings.MODEL_ID
        "doc_path": str,             # resolved absolute path
        "doc_name": str,             # file basename, for display/citation
        "client": str,               # the client this document belongs to
        "file_mtime": float,         # source file's mtime at index time
        "chunks": [{"text": str}, ...],
        "embeddings": np.ndarray,    # (N, dim) float32, L2-normalized
    }
(embeddings live in the .npz; everything else in the .json — see
utils/embedding_store.py.)

This intentionally mirrors the session embeddings sidecar
(services/search_service.py) field-for-field where the concepts
overlap (model_id, chunks, embeddings) so SearchService can merge both
into one flat matrix with minimal branching.

Security note: this used to be a single pickle file per document, and
pickle.loads() on an attacker-controlled file is arbitrary code
execution. doc_index/ lives inside the recordings directory, which is
synced across machines via Google Drive / OneDrive — so these files are
not purely local, and a compromised cloud account or bad sync could
have handed an attacker RCE with no local access at all. See
utils/embedding_store.py for the full rationale. Legacy `doc_*.pkl`
files from before this migration are never read, not even to migrate
them — index_folder() treats one as "not indexed" and rebuilds the
document's chunks and embeddings straight from the source file.
"""

from __future__ import annotations

import hashlib
import json
from html.parser import HTMLParser
from pathlib import Path
from typing import Callable, Dict, List, Optional, Tuple

import numpy as np

from core.embeddings import MODEL_ID
from utils.embedding_store import delete_payload, save_payload
from utils.logger import get_logger

logger = get_logger(__name__)

# Field data (2026-08-14 LMA gap analysis): a real client Knowledge
# Folder — SOWs, estimates, RFP responses, architecture decks — hit 87
# skipped files, nearly the user's whole corpus. Root causes: pypdf and
# python-docx were declared supported but never shipped, and .xlsx/.pptx/
# .html had no extractor at all despite openpyxl already being a hard
# dependency for the Excel export feature. This set (and the dispatch
# table below it, _EXTRACTORS) is now the actual source of truth for
# what's readable — see the consistency test in
# tests/test_document_service.py that fails if they ever drift apart
# again.
#
# Plain text: read straight off disk, no library at all. Every other
# extension must have a matching entry in _EXTRACTORS with a lazy
# import and a skip-reason path — never a hard dependency for the rest
# of the app to import this module.
#
# Field data (2026-08-19): a real client Knowledge Folder skipped 19
# files. Ten of those skips were correct (6 images, 3 archives, 1
# .drawio). The rest were coverage gaps, and the plain-text ones were
# gaps by oversight rather than by decision — the module was indexing
# .xlsx spreadsheets while reporting "unsupported file type: .csv",
# and CloudFormation .yaml templates are as plain-text as the .txt
# files right next to them.
#
# Everything here is read through _read_text_bounded(), which stops at
# the character cap instead of slurping the file — a stray 2 GB .log
# in a knowledge folder must not become a 2 GB allocation.
_PLAIN_TEXT_EXTENSIONS = {
    # prose / notes
    ".txt", ".md", ".rst", ".log",
    # config, IaC and data formats an SA's knowledge folder actually
    # holds: CloudFormation and Kubernetes YAML, Terraform, API
    # payloads, exported settings.
    ".yaml", ".yml", ".json", ".xml", ".ini", ".cfg", ".toml", ".tf",
    ".sql",
}

# 350 words is a slightly tighter window than the 400-word transcript
# chunking in core.embeddings — document prose tends to be denser
# (no filler speech) so a smaller window keeps chunks topically
# coherent for retrieval.
TARGET_CHUNK_WORDS = 350
OVERLAP_WORDS = 50

DOC_INDEX_SUBDIR = "doc_index"


def _doc_index_dir(recordings_dir) -> Path:
    return Path(recordings_dir) / DOC_INDEX_SUBDIR


def _doc_digest(doc_path: Path) -> str:
    return hashlib.sha1(str(doc_path.resolve()).encode("utf-8"), usedforsecurity=False).hexdigest()


def _doc_npz_path(recordings_dir, doc_path: Path) -> Path:
    """doc_<sha1(resolved path)>.npz — stable across reindex runs (same
    file always hashes to the same name), and collision-safe across
    clients/folders since the hash is over the full path."""
    return _doc_index_dir(recordings_dir) / f"doc_{_doc_digest(doc_path)}.npz"


def _doc_json_path(recordings_dir, doc_path: Path) -> Path:
    return _doc_index_dir(recordings_dir) / f"doc_{_doc_digest(doc_path)}.json"


def _doc_legacy_pkl_path(recordings_dir, doc_path: Path) -> Path:
    """Pre-migration pickle sidecar for this document. NEVER loaded —
    see module docstring and utils/embedding_store.py. Only ever checked
    for existence or unlinked."""
    return _doc_index_dir(recordings_dir) / f"doc_{_doc_digest(doc_path)}.pkl"


def _iter_candidate_files(folder: Path):
    """Every real file under `folder`, recursively, skipping any path
    with a hidden (dot-prefixed) component — .git, .DS_Store
    companions, sync-client sentinel dirs, etc. are tooling noise, never
    a document a user meant to index, so they're excluded entirely
    rather than reported as a skip.

    Deliberately NOT filtered by SUPPORTED_EXTENSIONS here: a .csv or
    .drawio the user dropped in the folder should show up in the
    reindex report as a skip with an "unsupported file type" (or "not a
    text document") reason (from extract_text), not vanish without a
    trace — the same silent-skip failure mode the module docstring
    calls out.
    """
    for path in folder.rglob("*"):
        if not path.is_file():
            continue
        try:
            rel_parts = path.relative_to(folder).parts
        except ValueError:
            continue
        if any(part.startswith(".") for part in rel_parts):
            continue
        yield path


# ── Text extraction ─────────────────────────────────────────────────

# Applied uniformly across every extractor, not just the newer ones —
# a single pathological source file (a hundred-sheet workbook, a
# thousand-page PDF) shouldn't be able to blow up chunk count/embedding
# time regardless of format. ~500k characters is roughly 85k words —
# generous for a real SOW/estimate/RFP response, a hard stop for
# anything absurd. Applied once, centrally, in extract_text() below
# rather than per-extractor so every format gets the same guarantee.
MAX_EXTRACTED_CHARS = 500_000

# Formats that are never going to hold extractable text — images and
# diagram/media/archive formats. Reported as "not a text document",
# never as a defect the user should go fix (see extract_text
# docstring). Values are the plural noun used in the skip message.
#
# .drawio is deliberately NOT here as an extraction target: modern
# draw.io files store each <diagram> as base64+deflate-compressed,
# URL-encoded XML (desktop-app default), while older/plain exports are
# uncompressed XML with mxCell value="..." labels — reliably telling
# the two apart and decoding the compressed variant needs a real test
# fixture from an actual draw.io export to get right, which we don't
# have here. That didn't "fall out cleanly," so .drawio is classified
# as a non-text diagram format below instead of guessing at a decoder.
NON_TEXT_EXTENSIONS: Dict[str, str] = {
    ".jpg": "images", ".jpeg": "images", ".png": "images", ".gif": "images",
    ".bmp": "images", ".tiff": "images", ".tif": "images", ".webp": "images",
    ".svg": "images", ".ico": "images", ".heic": "images",
    ".drawio": "diagrams", ".vsdx": "diagrams",
    ".mp3": "audio files", ".wav": "audio files", ".m4a": "audio files",
    ".mp4": "video files", ".mov": "video files",
    ".zip": "archives", ".rar": "archives", ".7z": "archives",
}


def extract_text(path: Path) -> Tuple[str, Optional[str]]:
    """Extract plain text from a supported document.

    Returns (text, skip_reason). Exactly one of the two is meaningful:
    on success `text` is non-empty and `skip_reason` is None; on any
    failure `text` is "" and `skip_reason` explains why in terms a user
    can act on:
      - a genuinely missing optional library ("pypdf isn't installed —
        this copy of the app is missing a component; reinstalling
        should fix it") — deliberately NOT "pip install X": inside a
        packaged desktop app with a bootstrapped venv the user has no
        pip to run, so that advice was actionable only for developers,
        misleading for everyone else.
      - a format that was never going to be text ("not a text document
        — images aren't indexed") — expected, not a problem, so it
        reads differently from a real failure.
      - an actual failure: corrupt/encrypted file, empty extraction,
        unsupported extension.

    Never raises for a bad input file — every failure mode observed in
    the field (missing optional dependency, corrupt/encrypted file,
    unsupported extension, empty extraction) is a returned reason, not
    an exception, so index_folder() can keep going and still report the
    skip instead of losing the whole reindex run to one bad file.
    """
    suffix = path.suffix.lower()
    try:
        if suffix in _PLAIN_TEXT_EXTENSIONS:
            text = _read_text_bounded(path)
            reason = None
        elif suffix in _EXTRACTORS:
            text, reason = _EXTRACTORS[suffix](path)
        elif suffix in NON_TEXT_EXTENSIONS:
            kind = NON_TEXT_EXTENSIONS[suffix]
            return "", f"not a text document — {kind} aren't indexed"
        else:
            return "", f"unsupported file type: {suffix or '(no extension)'}"
    except OSError as e:
        return "", f"could not read file: {e}"

    if reason:
        return "", reason

    text = (text or "").strip()
    if not text:
        return "", "no extractable text (empty or image-only document)"
    if len(text) > MAX_EXTRACTED_CHARS:
        text = text[:MAX_EXTRACTED_CHARS]
    return text, None


def _read_text_bounded(path: Path, limit: int = MAX_EXTRACTED_CHARS) -> str:
    """Read at most ``limit`` characters of a text file.

    extract_text() re-applies the cap on whatever comes back, so this
    exists purely to bound the ALLOCATION: ``read_text()`` on a 2 GB log
    file that someone dropped in a knowledge folder would materialise
    the whole thing in memory only to throw 99.9% of it away.
    """
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return f.read(limit + 1)


def _extract_csv(path: Path) -> Tuple[str, Optional[str]]:
    """One line of ``cell | cell | cell`` per row.

    Read through the ``csv`` module rather than splitting on commas so
    a quoted field containing a comma — or an embedded newline, which
    is common in exported estimates and RFP response matrices — stays
    in one piece instead of shattering the row.

    The output shape deliberately matches ``_extract_xlsx``'s
    (pipe-joined cells, blank cells and blank rows dropped): a CSV is
    the same data as a one-sheet workbook, so it should chunk the same
    way and retrieve the same way. Without this, the whole file would
    collapse toward one unreadable line and chunk boundaries would fall
    mid-record.

    Sniffs the delimiter so tab-separated files (.tsv, and the many
    .csv files that are actually tab- or semicolon-delimited) do not
    come back as one column.
    """
    import csv  # noqa: PLC0415  (stdlib, but keep the import local)

    try:
        sample = _read_text_bounded(path, 8192)
    except OSError as e:
        return "", f"could not read file: {e}"
    if not sample.strip():
        return "", None  # extract_text() turns this into "no extractable text"

    try:
        dialect = csv.Sniffer().sniff(sample, delimiters=",;\t|")
    except csv.Error:
        # A single-column file gives the sniffer nothing to go on. That
        # is not an error — comma is the right default.
        dialect = csv.excel

    lines: List[str] = []
    total = 0
    try:
        with open(path, "r", encoding="utf-8", errors="replace",
                  newline="") as f:
            for row in csv.reader(f, dialect):
                # Collapse whitespace INSIDE each cell. A quoted field
                # may legally contain newlines (common in RFP response
                # matrices and estimate notes), and leaving them in
                # would split one record across several output lines —
                # defeating the point of doing this row-wise at all.
                cells = [" ".join(str(c).split()) for c in row]
                cells = [c for c in cells if c]
                if not cells:
                    continue
                line = " | ".join(cells)
                lines.append(line)
                total += len(line) + 1
                # Same bail-out as _extract_xlsx: stop once past the cap
                # rather than reading a huge file only to truncate it.
                if total > MAX_EXTRACTED_CHARS:
                    break
    except (OSError, csv.Error) as e:
        return "", f"corrupt or unreadable csv: {e}"

    return "\n".join(lines), None


# ── legacy .doc ──────────────────────────────────────────────────────
#
# Field data (2026-08-19): 5 of the 19 skipped files in a real client
# Knowledge Folder were `.doc` — and they were the valuable ones (TCO/ROI
# models, a competitive battle card, an architecture document).
#
# `.doc` is not one format. The extension is attached to at least four
# different things in the wild, and which one you have decides whether
# it is trivial or impossible to read:
#
#   1. RTF with a .doc extension. Extremely common — "Save as .doc"
#      from non-Word tools, and most mail/CRM exports, produce this.
#   2. HTML with a .doc extension. Word's own "Save as Web Page", and
#      almost every server-side "export to Word" feature, produce this.
#   3. Plain text with a .doc extension.
#   4. A genuine Word 97-2003 OLE2 compound file.
#
# So the extractor sniffs the magic bytes first and routes; only case 4
# needs any real work. python-docx cannot read ANY of these — it is a
# zip reader, and none of the four is a zip.
_OLE2_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"


def _sniff_document_kind(path: Path) -> str:
    """What a `.doc` file actually is: ole2 / rtf / html / text / binary."""
    try:
        with open(path, "rb") as f:
            head = f.read(4096)
    except OSError:
        return "binary"

    if head.startswith(_OLE2_MAGIC):
        return "ole2"
    stripped = head.lstrip(b"\xef\xbb\xbf \t\r\n")
    if stripped.startswith(b"{\\rt"):
        return "rtf"
    lowered = stripped[:1024].lower()
    if (lowered.startswith(b"<!doctype html") or lowered.startswith(b"<html")
            or b"<html" in lowered or b"<body" in lowered):
        return "html"
    # Everything left: call it text if it decodes cleanly and is not
    # peppered with the NULs and control bytes that mark a binary blob.
    if not head:
        return "text"
    controls = sum(
        1 for b in head
        if b < 0x09 or (0x0E <= b < 0x20 and b not in (0x1B,))
    )
    if controls / len(head) < 0.02:
        return "text"
    return "binary"


def _rtf_to_text(raw: str) -> str:
    """Strip RTF markup down to readable prose.

    Deliberately small and approximate rather than a real RTF parser:
    this feeds semantic search, where approximate text beats no text.
    It drops the header groups that hold no prose (font/colour/style
    tables and any ``{\\*\\...}`` destination), resolves ``\\'hh`` byte
    escapes and the handful of control words that mean whitespace, and
    throws the remaining control words away.
    """
    out: List[str] = []
    i = 0
    n = len(raw)
    depth = 0
    skip_to_depth: Optional[int] = None

    while i < n:
        ch = raw[i]
        if ch == "{":
            depth += 1
            # Groups that never contain document prose.
            ahead = raw[i + 1:i + 12].lower()
            if skip_to_depth is None and (
                ahead.startswith("\\*")
                or ahead.startswith("\\fonttbl")
                or ahead.startswith("\\colortbl")
                or ahead.startswith("\\stylesheet")
                or ahead.startswith("\\info")
                or ahead.startswith("\\pict")
            ):
                skip_to_depth = depth
            i += 1
            continue
        if ch == "}":
            if skip_to_depth is not None and depth == skip_to_depth:
                skip_to_depth = None
            depth -= 1
            i += 1
            continue
        if ch == "\\":
            # \'hh — a single byte in the document's codepage.
            if raw[i + 1:i + 2] == "'":
                hexpair = raw[i + 2:i + 4]
                i += 4
                if skip_to_depth is None:
                    try:
                        out.append(bytes([int(hexpair, 16)]).decode(
                            "cp1252", errors="replace"))
                    except ValueError:
                        pass
                continue
            # An escaped literal character.
            if not raw[i + 1:i + 2].isalpha():
                if skip_to_depth is None and raw[i + 1:i + 2]:
                    out.append(raw[i + 1])
                i += 2
                continue
            # A control word, optionally with a numeric parameter and
            # one optional trailing space that belongs to the word.
            j = i + 1
            while j < n and raw[j].isalpha():
                j += 1
            word = raw[i + 1:j]
            if j < n and (raw[j] == "-" or raw[j].isdigit()):
                j += 1
                while j < n and raw[j].isdigit():
                    j += 1
            if j < n and raw[j] == " ":
                j += 1
            if skip_to_depth is None and word in (
                    "par", "line", "tab", "cell", "row", "sect", "page"):
                out.append("\n" if word != "tab" else "\t")
            i = j
            continue
        if skip_to_depth is None and ch not in "\r\n":
            out.append(ch)
        i += 1

    text = "".join(out)
    # RTF line breaks are the control words above; raw newlines in the
    # source are formatting noise. Collapse the runs we produced.
    lines = [ln.strip() for ln in text.splitlines()]
    return "\n".join(ln for ln in lines if ln)


# Offsets into the WordDocument stream's FIB, from [MS-DOC]. The FIB is
# FibBase (0x00-0x1F), csw (0x20), fibRgW (0x22-0x3D), cslw (0x3E),
# fibRgLw97 (0x40-0x97), cbRgFcLcb (0x98), then the FibRgFcLcb97 blob
# from 0x9A. fcClx is pair 33 of that blob: 0x9A + 33*8 = 0x1A2.
_FIB_FLAGS_OFFSET = 0x0A
_FIB_WHICH_TABLE_STREAM_BIT = 0x0200
_FIB_FCCLX_OFFSET = 0x01A2


def _extract_doc_ole2(path: Path) -> Tuple[str, Optional[str]]:
    """Text out of a genuine Word 97-2003 binary document.

    Walks the piece table exactly as [MS-DOC] specifies rather than
    scraping printable runs out of the stream: the WordDocument stream
    interleaves document text with field codes, style names and binary
    structures, so a "pull out the readable-looking bytes" approach
    yields prose salted with garbage — which is worse than nothing for
    a semantic index, because the garbage gets embedded and retrieved
    alongside real content.

    The walk is: read the FIB to find which table stream is live and
    where the Clx sits, skip the RgPrc grpprl entries, then read the
    PlcPcd. Each piece's fc has bit 30 set when the text is 8-bit
    cp1252 ("compressed") at fc/2, and clear when it is UTF-16LE at fc.

    Validated against a real Word 97-2003 document, not a synthetic
    one — see tests/fixtures/README.md.

    olefile is the only dependency, and it is deliberately the ONLY one
    considered acceptable here: pure Python, no transitive
    dependencies, ~115 KB. LibreOffice-headless and textract were both
    ruled out — neither can exist inside a packaged app.
    """
    try:
        import olefile  # noqa: PLC0415  (deliberately lazy, as above)
    except ImportError:
        return "", _missing_dependency_reason("olefile")

    import struct  # noqa: PLC0415

    try:
        ole = olefile.OleFileIO(str(path))
    except Exception as e:
        return "", f"corrupt or unreadable doc: {e}"

    try:
        if not ole.exists("WordDocument"):
            # An OLE2 file that is not a Word document at all — an old
            # .xls or .ppt renamed, most often. Say so precisely.
            return "", ("not a Word document — this is an OLE2 file with "
                        "no WordDocument stream")
        wd = ole.openstream("WordDocument").read()
        if len(wd) < _FIB_FCCLX_OFFSET + 8:
            return "", "corrupt or unreadable doc: truncated FIB"

        flags = struct.unpack_from("<H", wd, _FIB_FLAGS_OFFSET)[0]
        table_name = ("1Table" if flags & _FIB_WHICH_TABLE_STREAM_BIT
                      else "0Table")
        if not ole.exists(table_name):
            return "", (f"corrupt or unreadable doc: missing {table_name} "
                        f"stream")
        table = ole.openstream(table_name).read()

        fc_clx, lcb_clx = struct.unpack_from("<II", wd, _FIB_FCCLX_OFFSET)
        clx = table[fc_clx:fc_clx + lcb_clx]
        if not clx:
            return "", "corrupt or unreadable doc: empty piece table"

        # Skip any RgPrc (clxtGrpprl == 0x01) entries preceding the Pcdt.
        i = 0
        while i < len(clx) and clx[i] == 0x01:
            cb = struct.unpack_from("<H", clx, i + 1)[0]
            i += 3 + cb
        if i >= len(clx) or clx[i] != 0x02:
            return "", "corrupt or unreadable doc: no piece table found"

        lcb = struct.unpack_from("<I", clx, i + 1)[0]
        plc = clx[i + 5:i + 5 + lcb]
        piece_count = (len(plc) - 4) // 12
        if piece_count <= 0:
            return "", "corrupt or unreadable doc: empty piece table"

        cps = struct.unpack_from(f"<{piece_count + 1}I", plc, 0)
        parts: List[str] = []
        total = 0
        for k in range(piece_count):
            off = 4 * (piece_count + 1) + k * 8
            fc = struct.unpack_from("<I", plc, off + 2)[0]
            compressed = bool(fc & 0x40000000)
            base = fc & 0x3FFFFFFF
            nchars = cps[k + 1] - cps[k]
            if nchars <= 0:
                continue
            if compressed:
                chunk = wd[base // 2: base // 2 + nchars].decode(
                    "cp1252", errors="replace")
            else:
                chunk = wd[base: base + nchars * 2].decode(
                    "utf-16-le", errors="replace")
            parts.append(chunk)
            total += len(chunk)
            if total > MAX_EXTRACTED_CHARS:
                break
    except Exception as e:
        return "", f"corrupt or unreadable doc: {e}"
    finally:
        try:
            ole.close()
        except Exception:
            pass

    return _clean_doc_text("".join(parts)), None


# Word's binary text stream carries in-band markers alongside the prose:
# 0x07 ends a table cell/row, 0x0B is a hard line break, 0x0C a page
# break, 0x13/0x14/0x15 bracket field codes, and 0x1E/0x1F are
# non-breaking/optional hyphens.
_DOC_FIELD_START = "\x13"
_DOC_FIELD_SEPARATOR = "\x14"
_DOC_FIELD_END = "\x15"


def _clean_doc_text(raw: str) -> str:
    """Turn the raw piece-table text into readable prose.

    Drops field CODES while keeping field RESULTS — a HYPERLINK or REF
    field stores its instruction between 0x13 and 0x14 and the text the
    reader actually sees between 0x14 and 0x15, so keeping the former
    would index "HYPERLINK \\l _Toc12345" instead of the heading it
    points at.
    """
    out: List[str] = []
    in_field_code = False
    for ch in raw:
        if ch == _DOC_FIELD_START:
            in_field_code = True
            continue
        if ch == _DOC_FIELD_SEPARATOR:
            in_field_code = False
            continue
        if ch == _DOC_FIELD_END:
            in_field_code = False
            continue
        if in_field_code:
            continue
        if ch in ("\r", "\x0b", "\x0c", "\x07"):
            out.append("\n")
        elif ch in ("\x1e", "\x1f"):
            out.append("-")
        elif ch == "\t":
            out.append("\t")
        elif ch < " " and ch != "\n":
            continue
        else:
            out.append(ch)
    lines = [ln.strip() for ln in "".join(out).splitlines()]
    return "\n".join(ln for ln in lines if ln)


def _extract_rtf(path: Path) -> Tuple[str, Optional[str]]:
    try:
        raw = _read_text_bounded(path, MAX_EXTRACTED_CHARS * 4)
    except OSError as e:
        return "", f"could not read file: {e}"
    return _rtf_to_text(raw), None


def _extract_doc(path: Path) -> Tuple[str, Optional[str]]:
    """`.doc` — sniff what the file really is, then route.

    See the block comment above ``_OLE2_MAGIC`` for why sniffing comes
    first: the extension says almost nothing about the format.
    """
    kind = _sniff_document_kind(path)
    if kind == "ole2":
        return _extract_doc_ole2(path)
    if kind == "rtf":
        return _extract_rtf(path)
    if kind == "html":
        return _extract_html(path)
    if kind == "text":
        try:
            return _read_text_bounded(path), None
        except OSError as e:
            return "", f"could not read file: {e}"
    return "", ("unrecognised .doc format — not an OLE2 Word document, "
                "RTF, HTML or text")


def _missing_dependency_reason(package: str) -> str:
    """Shared phrasing for every lazy-imported extractor library that
    turns out not to be installed. NOT "pip install X" — see
    extract_text docstring for why that's the wrong advice inside a
    packaged app."""
    return (
        f"{package} isn't installed — this copy of the app is missing "
        "a component; reinstalling should fix it"
    )


def _extract_pdf(path: Path) -> Tuple[str, Optional[str]]:
    try:
        import pypdf  # noqa: PLC0415  (deliberately lazy — see module docstring)
    except ImportError:
        return "", _missing_dependency_reason("pypdf")

    try:
        reader = pypdf.PdfReader(str(path))
    except Exception as e:
        return "", f"corrupt or unreadable PDF: {e}"

    if reader.is_encrypted:
        try:
            # Some "encrypted" PDFs are only permission-restricted with
            # no real open password — an empty-password decrypt often
            # succeeds for those. A genuinely password-protected file
            # fails here and we report it rather than crash.
            if not reader.decrypt(""):
                return "", "encrypted PDF (password protected) — could not open"
        except Exception:
            return "", "encrypted PDF (password protected) — could not open"

    try:
        pages_text = [page.extract_text() or "" for page in reader.pages]
    except Exception as e:
        return "", f"corrupt or unreadable PDF: {e}"
    return "\n".join(pages_text), None


def _extract_docx(path: Path) -> Tuple[str, Optional[str]]:
    try:
        import docx  # noqa: PLC0415  (python-docx; deliberately lazy)
    except ImportError:
        return "", _missing_dependency_reason("python-docx")

    try:
        document = docx.Document(str(path))
    except Exception as e:
        return "", f"corrupt or unreadable docx: {e}"

    paragraphs = [p.text for p in document.paragraphs]
    return "\n".join(paragraphs), None


def _extract_xlsx(path: Path) -> Tuple[str, Optional[str]]:
    """Sheet-by-sheet cell text, useful for semantic search: each
    sheet's name as a header, then its rows as pipe-joined cell text.
    Blank cells and blank rows are dropped so the extracted text isn't
    mostly whitespace for sparse workbooks.

    read_only + data_only load: read_only streams rows without loading
    charts/images into memory at all (they're simply never visited),
    and data_only returns each formula cell's last-calculated value
    instead of the formula string — "=SUM(B2:B40)" is useless for
    search, "142500" is not.
    """
    try:
        import openpyxl  # noqa: PLC0415  (deliberately lazy — see module docstring)
    except ImportError:
        return "", _missing_dependency_reason("openpyxl")

    try:
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    except Exception as e:
        return "", f"corrupt or unreadable xlsx: {e}"

    parts: List[str] = []
    try:
        for ws in wb.worksheets:
            sheet_lines = [f"## {ws.title}"]
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None and str(c).strip()]
                if cells:
                    sheet_lines.append(" | ".join(cells))
            if len(sheet_lines) > 1:
                parts.append("\n".join(sheet_lines))
            # Bail out once well past the cap rather than fully reading
            # (and holding in memory) every sheet of a huge workbook
            # only to truncate it afterward — extract_text() re-applies
            # the exact cap on the joined result either way.
            if sum(len(p) for p in parts) > MAX_EXTRACTED_CHARS:
                break
    except Exception as e:
        return "", f"corrupt or unreadable xlsx: {e}"
    finally:
        wb.close()

    return "\n\n".join(parts), None


def _extract_pptx(path: Path) -> Tuple[str, Optional[str]]:
    """Per-slide shape text plus speaker notes. python-pptx pulls in
    Pillow and lxml transitively — a heavier addition than pypdf or
    python-docx (see requirements files for the measured install
    size) — so this stays behind the same lazy-import/skip-reason
    pattern as every other optional extractor: a missing library here
    degrades to a clear, counted skip, never an ImportError at module
    load time for the rest of the app.
    """
    try:
        from pptx import Presentation  # noqa: PLC0415  (python-pptx; deliberately lazy)
    except ImportError:
        return "", _missing_dependency_reason("python-pptx")

    try:
        prs = Presentation(str(path))
    except Exception as e:
        return "", f"corrupt or unreadable pptx: {e}"

    parts: List[str] = []
    try:
        for i, slide in enumerate(prs.slides, start=1):
            slide_lines = [f"## Slide {i}"]
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    t = (shape.text_frame.text or "").strip()
                    if t:
                        slide_lines.append(t)
                elif getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            slide_lines.append(" | ".join(cells))
            if getattr(slide, "has_notes_slide", False):
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes:
                    slide_lines.append(f"Speaker notes: {notes}")
            if len(slide_lines) > 1:
                parts.append("\n".join(slide_lines))
    except Exception as e:
        return "", f"corrupt or unreadable pptx: {e}"

    return "\n\n".join(parts), None


class _HTMLTextExtractor(HTMLParser):
    """Minimal tag-stripping text extractor — stdlib only, no new
    dependency. Drops <script>/<style> contents (never document
    prose); keeps everything else, one text run per line."""

    def __init__(self) -> None:
        super().__init__()
        self._parts: List[str] = []
        self._skip_depth = 0

    def handle_starttag(self, tag: str, attrs) -> None:  # noqa: ARG002
        if tag in ("script", "style"):
            self._skip_depth += 1

    def handle_endtag(self, tag: str) -> None:
        if tag in ("script", "style") and self._skip_depth > 0:
            self._skip_depth -= 1

    def handle_data(self, data: str) -> None:
        if self._skip_depth == 0 and data.strip():
            self._parts.append(data.strip())

    def get_text(self) -> str:
        return "\n".join(self._parts)


def _extract_html(path: Path) -> Tuple[str, Optional[str]]:
    raw = path.read_text(encoding="utf-8", errors="replace")
    parser = _HTMLTextExtractor()
    try:
        parser.feed(raw)
        parser.close()
    except Exception as e:
        return "", f"corrupt or unreadable html: {e}"
    return parser.get_text(), None


# Dispatch table: the single source of truth for "which extensions can
# actually be read" — SUPPORTED_EXTENSIONS below is derived from this
# plus the inline-handled plain-text extensions, so the two can never
# silently drift apart the way "unsupported file type: .xlsx" (despite
# openpyxl being bundled) drifted from the real capability of this
# module. See test_supported_extensions_and_dispatch_agree.
_EXTRACTORS: Dict[str, Callable[[Path], Tuple[str, Optional[str]]]] = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".xlsx": _extract_xlsx,
    ".pptx": _extract_pptx,
    ".html": _extract_html,
    ".htm": _extract_html,
    # Row-aware, so records stay coherent across chunk boundaries and
    # the output matches _extract_xlsx's shape.
    ".csv": _extract_csv,
    ".tsv": _extract_csv,
    # Legacy Word. Sniffs the magic bytes and routes — see the block
    # comment above _OLE2_MAGIC; the extension covers four unrelated
    # formats.
    ".doc": _extract_doc,
    # Falls out of the .doc work for free: an actual .rtf gets the same
    # de-markup path a mislabelled one does.
    ".rtf": _extract_rtf,
}

SUPPORTED_EXTENSIONS = _PLAIN_TEXT_EXTENSIONS | set(_EXTRACTORS)


# ── Chunking ─────────────────────────────────────────────────────────

def chunk_text(
    text: str,
    target_words: int = TARGET_CHUNK_WORDS,
    overlap_words: int = OVERLAP_WORDS,
) -> List[str]:
    """Split `text` into ~target_words chunks with overlap_words shared
    between consecutive chunks, so a topic straddling a chunk boundary
    still has a full match in at least one chunk.

    Pure function — no file or model access — so it's trivially
    unit-testable. Short text (<= target_words) comes back as a single
    chunk; empty/whitespace-only text comes back as an empty list.
    """
    words = text.split()
    if not words:
        return []
    if len(words) <= target_words:
        return [" ".join(words)]

    # Guarantee forward progress even if a caller passes
    # overlap_words >= target_words.
    step = max(1, target_words - overlap_words)
    chunks: List[str] = []
    start = 0
    n = len(words)
    while start < n:
        end = min(start + target_words, n)
        chunks.append(" ".join(words[start:end]))
        if end >= n:
            break
        start += step
    return chunks


# ── Indexing ─────────────────────────────────────────────────────────

def _normalize_rows(matrix: np.ndarray) -> np.ndarray:
    """L2-normalize each row, defensively. embed_fn is caller-supplied
    (tests inject a fake) — normalizing here, not just trusting the
    caller, is what keeps the persisted sidecar matching the session
    embeddings sidecar's "L2-normalized float32" convention exactly, so
    SearchService can dot-product session and document vectors in the
    same matrix without a per-source special case."""
    matrix = np.asarray(matrix, dtype=np.float32)
    if matrix.size == 0:
        return matrix
    norms = np.linalg.norm(matrix, axis=1, keepdims=True)
    norms[norms == 0] = 1.0
    return (matrix / norms).astype(np.float32)


def index_folder(
    folder,
    client: str,
    embed_fn: Callable[[List[str]], np.ndarray],
    recordings_dir,
) -> dict:
    """Extract, chunk, and embed every supported document under
    `folder`, persisting one .npz/.json sidecar pair per document under
    <recordings_dir>/doc_index/.

    embed_fn is injected (callable: List[str] -> ndarray) so callers
    that can't/won't load sentence-transformers (all backend tests)
    can pass a fake embedder — this module never imports
    sentence-transformers itself.

    Re-embedding is skipped ("unchanged") when the existing sidecar is
    newer than the source file's mtime — a reindex of an unchanged
    folder should be near-instant. A document reachable only through a
    legacy .pkl sidecar (pre-migration, never loaded — see module
    docstring) is NOT "unchanged": it has no readable .npz/.json pair,
    so it's rebuilt from the source file same as an unindexed document,
    and the stale .pkl is removed once the rebuild succeeds.

    Returns {"indexed": int, "unchanged": int,
             "skipped": [{"file": str, "reason": str, "expected": bool}, ...],
             "total_chunks": int} — total_chunks counts every chunk
    currently on disk for this run (both freshly indexed and unchanged
    documents), so the report reflects the whole current index rather
    than only this run's work.

    Each skip's "expected" flag distinguishes a format that was never
    going to be a text document (images, diagrams, media, archives —
    see NON_TEXT_EXTENSIONS) from a genuine failure (missing library,
    corrupt file, unsupported extension, empty extraction). The
    frontend uses this to avoid making the former read like a defect
    the user should go fix.
    """
    folder = Path(folder)
    doc_dir = _doc_index_dir(recordings_dir)
    report = {"indexed": 0, "unchanged": 0, "skipped": [], "total_chunks": 0}

    if not folder.is_dir():
        report["skipped"].append(
            {"file": str(folder), "reason": "folder does not exist",
             "expected": False})
        _emit_index_event(report)
        return report

    doc_dir.mkdir(parents=True, exist_ok=True)

    for path in _iter_candidate_files(folder):
        npz_path = _doc_npz_path(recordings_dir, path)
        json_path = _doc_json_path(recordings_dir, path)
        legacy_pkl_path = _doc_legacy_pkl_path(recordings_dir, path)
        try:
            file_mtime = path.stat().st_mtime
        except OSError as e:
            report["skipped"].append(
                {"file": str(path), "reason": f"could not stat file: {e}",
                 "expected": False})
            continue

        if npz_path.exists() and json_path.exists():
            try:
                sidecar_mtime = min(
                    npz_path.stat().st_mtime, json_path.stat().st_mtime)
                sidecar_is_current = sidecar_mtime >= file_mtime
            except OSError:
                sidecar_is_current = False
            if sidecar_is_current:
                report["unchanged"] += 1
                try:
                    meta = json.loads(json_path.read_text(encoding="utf-8"))
                    report["total_chunks"] += len(meta.get("chunks") or [])
                except (OSError, ValueError) as e:
                    logger.warning(
                        f"Unchanged sidecar {json_path.name} unreadable "
                        f"for chunk count ({e}); count omitted, file kept.")
                continue

        text, reason = extract_text(path)
        if reason:
            report["skipped"].append({
                "file": str(path), "reason": reason,
                "expected": path.suffix.lower() in NON_TEXT_EXTENSIONS,
            })
            continue

        chunks = chunk_text(text)
        if not chunks:
            report["skipped"].append(
                {"file": str(path),
                 "reason": "no chunks produced from extracted text",
                 "expected": False})
            continue

        try:
            embeddings = embed_fn(chunks)
        except Exception as e:
            report["skipped"].append(
                {"file": str(path), "reason": f"embedding failed: {e}",
                 "expected": False})
            continue

        embeddings = _normalize_rows(embeddings)
        if embeddings.shape[0] != len(chunks):
            report["skipped"].append(
                {"file": str(path),
                 "reason": (
                     f"embed_fn returned {embeddings.shape[0]} vectors "
                     f"for {len(chunks)} chunks"),
                 "expected": False})
            continue

        meta = {
            "model_id": MODEL_ID,
            "doc_path": str(path.resolve()),
            "doc_name": path.name,
            "client": client,
            "file_mtime": file_mtime,
            "chunks": [{"text": t} for t in chunks],
        }
        save_payload(npz_path, json_path, embeddings=embeddings, meta=meta)
        if legacy_pkl_path.exists():
            # Freshly rebuilt from source — the legacy pickle for this
            # document is now redundant and, left in place, is a
            # standing load hazard if any future code path regressed to
            # reading it. Clean it up.
            delete_payload(legacy_pkl_path)
            logger.info(f"Removed legacy pickle sidecar {legacy_pkl_path.name}")

        report["indexed"] += 1
        report["total_chunks"] += len(chunks)
        logger.info(
            f"Indexed document {path.name}: {len(chunks)} chunks "
            f"({client}) -> {npz_path.name}")

    _emit_index_event(report)
    return report


# Skip reasons are written for humans and interpolate the offending
# file's path and the library's exception message straight into the
# string. Neither may ever reach events.jsonl, so skips are reduced to
# this fixed set of categories before they are counted. The prose
# reasons stay exactly as they are for the UI and backend.log.
_SKIP_CATEGORY_MARKERS = (
    ("not a text document", "not_a_text_document"),
    ("unsupported file type", "unsupported_file_type"),
    ("no extractable text", "no_extractable_text"),
    ("no chunks produced", "no_chunks_produced"),
    ("isn't installed", "missing_dependency"),
    ("encrypted", "encrypted"),
    ("corrupt or unreadable", "corrupt_or_unreadable"),
    ("could not stat file", "could_not_stat"),
    ("could not read file", "could_not_read"),
    ("embedding failed", "embedding_failed"),
    ("returned", "embedder_count_mismatch"),
    ("folder does not exist", "folder_missing"),
)


def skip_reason_category(reason: str) -> str:
    """Map a human skip reason onto a stable, path-free enum token."""
    low = (reason or "").lower()
    for marker, category in _SKIP_CATEGORY_MARKERS:
        if marker in low:
            return category
    return "other"


def _emit_index_event(report: dict) -> None:
    """One ``documents.indexed`` event per reindex run.

    Counts and file EXTENSIONS only. The client name is a real customer
    name and the skipped files' paths are the user's directory tree —
    neither goes in a file meant to be attached to a bug report, so
    neither is passed here. The v2.28.0 field report ("87 skipped files
    out of a whole corpus") is answerable from the by-reason and
    by-extension histograms alone, which is the point.
    """
    try:
        from utils import events
        by_reason: Dict[str, int] = {}
        by_extension: Dict[str, int] = {}
        expected_skips = 0
        for s in report.get("skipped") or []:
            cat = skip_reason_category(str(s.get("reason") or ""))
            by_reason[cat] = by_reason.get(cat, 0) + 1
            if s.get("expected"):
                expected_skips += 1
            suffix = Path(str(s.get("file") or "")).suffix.lower()
            key = (suffix.lstrip(".") or "none")
            if key.isalnum() and len(key) <= 12:
                by_extension[key] = by_extension.get(key, 0) + 1
        events.emit(
            events.DOCUMENTS_INDEXED,
            indexed=int(report.get("indexed") or 0),
            unchanged=int(report.get("unchanged") or 0),
            skipped=len(report.get("skipped") or []),
            expected_skips=expected_skips,
            total_chunks=int(report.get("total_chunks") or 0),
            skipped_by_reason=by_reason,
            skipped_by_extension=by_extension,
        )
    except Exception:
        pass


def remove_stale(folder, client: str, recordings_dir) -> int:
    """Delete doc_index sidecars for `client` whose source file no
    longer exists on disk — the user deleted or moved a document out of
    the knowledge folder since the last reindex.

    Also opportunistically removes EVERY leftover legacy `doc_*.pkl`
    sidecar in doc_index/, regardless of client: those files are never
    loaded (see module docstring) and there's no safe way to read which
    client or source path one belongs to without unpickling it. They're
    pure clutter/hazard at this point — deleting one just means the
    owning document gets rebuilt from source next time its client's
    folder is reindexed, which is always safe (indexes rebuild lazily).

    `folder` isn't used to further scope the scan (client already
    disambiguates within recordings_dir) but is kept in the signature
    for symmetry with index_folder() and because a future version may
    want to distinguish "moved to a different knowledge folder" from
    "deleted".
    """
    del folder  # not needed for staleness (see docstring); kept for API symmetry
    doc_dir = _doc_index_dir(recordings_dir)
    if not doc_dir.is_dir():
        return 0

    for legacy in doc_dir.glob("doc_*.pkl"):
        logger.warning(
            f"Ignoring legacy pickle sidecar {legacy.name} "
            f"(never loaded — removing; owning document will be "
            f"rebuilt from source on next reindex)")
        try:
            legacy.unlink()
        except OSError as e:
            logger.warning(f"Could not delete {legacy}: {e}")

    removed = 0
    for json_path in doc_dir.glob("doc_*.json"):
        try:
            payload = json.loads(json_path.read_text(encoding="utf-8"))
        except (OSError, ValueError) as e:
            logger.warning(f"Could not read {json_path.name}: {e}")
            continue
        if (payload.get("client") or "") != client:
            continue
        doc_path_str = payload.get("doc_path") or ""
        if doc_path_str and Path(doc_path_str).exists():
            continue
        npz_path = json_path.with_suffix(".npz")
        if delete_payload(npz_path, json_path):
            removed += 1
            logger.info(
                f"Removed stale doc index {json_path.stem} "
                f"(source gone: {doc_path_str})")
    return removed
